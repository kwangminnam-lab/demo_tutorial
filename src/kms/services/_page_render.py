"""페이지 이미지 렌더 + 변경 사항 색칠 — 포맷 무관 PNG data URL 생성.

흐름:
  1. 입력 확장자가 지원 목록(`_RENDERABLE_EXTS`)이면 처리, 아니면 빈 dict.
  2. PDF는 바로 pymupdf로 페이지별 픽맵 → PNG.
  3. 비PDF (DOCX/PPTX/XLSX/HTML/TXT/MD)는 **LibreOffice headless**로 PDF 변환 →
     같은 pymupdf 경로로 합류. 변환 PDF는 임시 파일로 두고 렌더 후 즉시 삭제.
  4. **하이라이트** — 호출자가 `[(text, kind), ...]` 목록을 넘기면 각 페이지에서
     `page.search_for(text)`로 위치를 찾아 형광 annotation을 추가. 픽맵 렌더에
     자동 포함되어 PNG 결과에 변경 부분이 색으로 표시된다.
        kind: "add" → 초록, "delete" → 빨강, "change" → 노랑.
  5. 모든 단계의 실패(SDK 부재·LibreOffice 부재·변환 타임아웃·손상 파일·개별
     페이지 렌더 실패·텍스트 매칭 실패)는 WARNING/DEBUG 로그 + 빈 dict 또는
     부분 결과로 graceful degrade한다. 비교 본체(텍스트 diff)는 항상 진행한다.

설계 결정:
  - 단일 렌더 경로로 통일 — 비PDF도 결국 PDF로 정규화한 뒤 pymupdf. 포맷별
    분기·코드 중복 제거. PDF 페이지 개념이 가장 일반적이라 사용자 경험이 일관.
  - LibreOffice는 사내 운영 환경(컨테이너 이미지)에 포함시키는 비교적 표준
    경로. macOS·Linux 모두 동작. `soffice --convert-to pdf`.
  - 하이라이트는 PDF annotation 레이어로 추가 — 픽맵 렌더 시 자동 포함, 별도
    합성 코드 X. 색은 형광펜 stroke + opacity 0.35.
  - 캐시 키는 (원본 path, mtime_ns, highlights tuple) — 변환 PDF mtime은 무관.
  - 메모리 only 캐시 (디스크 없음 — ADR-007 연장, 사내 자료 잔존 위험 회피).
"""

from __future__ import annotations

import base64
import logging
import os
import re
import shutil
import subprocess
import tempfile
from functools import lru_cache
from pathlib import Path
from typing import Literal

logger = logging.getLogger(__name__)

#: 하이라이트 종류 — text diff 연산과 일치.
HighlightKind = Literal["add", "delete", "change"]

#: Highlight entry: (라인 전체, kind, 변경 단어 튜플).
#: - kind=="change" + changed_words 채움 → 라인 매치 위치 안 변경 단어만 색칠.
#: - kind=="add"/"delete" 또는 changed_words=() → 라인 전체 색칠.
Highlight = tuple[str, HighlightKind, tuple[str, ...]]

#: kind별 형광 색 (R, G, B) 0~1. 사용자 인지 색 규약 — 추가 초록, 삭제 빨강, 변경 노랑.
_HIGHLIGHT_COLORS: dict[str, tuple[float, float, float]] = {
    "add": (0.16, 0.78, 0.35),
    "delete": (0.86, 0.15, 0.15),
    "change": (0.96, 0.75, 0.20),
}

#: 형광 투명도 — 본문 글자 가독성 유지.
_HIGHLIGHT_OPACITY = 0.35

#: 한 텍스트당 PDF 검색 prefix 상한 — PDF 줄바꿈으로 매칭 실패 회피.
_SEARCH_PREFIX_MAX = 60

#: full-line 검색이 빗나갈 때 점진적으로 줄이는 prefix의 최소 길이.
#: 이보다 짧은 prefix는 흔한 단어와 우연 매칭(false positive) 위험이 커 검색하지 않는다.
_SEARCH_MIN_PREFIX = 10

#: 검색 대상에서 제외할 마커 라인 프리픽스(추출기가 박는 이미지·표 마커 등).
_MARKER_PREFIXES = ("[IMAGE", "[TABLE")

#: 추출기(markdown)가 괄호·구두점 둘레에 끼워 넣는 공백 인공물 보정용 문자군.
#: 예: 추출 '[ 런웨이 요금정책 ]' ↔ PDF 원문 '[런웨이 요금정책]'. 여는 괄호 뒤·
#: 닫는 괄호 앞·종결 구두점 앞의 공백만 제거해 search_for 매칭률을 높인다(단어
#: 사이 공백은 보존 — 과교정 방지).
_OPEN_PUNCT = "([{<「『【（《"
_CLOSE_PUNCT = ")]}>」』】）》"
_TERM_PUNCT = ",.·:;!?…"
_RE_SPACE_AFTER_OPEN = re.compile(r"([" + re.escape(_OPEN_PUNCT) + r"])\s+")
_RE_SPACE_BEFORE_CLOSE = re.compile(r"\s+([" + re.escape(_CLOSE_PUNCT + _TERM_PUNCT) + r"])")

#: 페이지 단위 렌더 대상 확장자 — 모든 지원 포맷 + 스캔 이미지.
_RENDERABLE_EXTS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xlsm",
    ".html",
    ".htm",
    ".txt",
    ".md",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
    ".gif",
}

_PDF_EXT = ".pdf"

#: 이미지(스캔)는 pymupdf로 PDF 변환해 렌더한다(LibreOffice 불요).
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif"}

#: 한 문서당 렌더 페이지 상한 — 메모리·payload 보호.
_DEFAULT_MAX_PAGES = 50

#: 픽맵 줌 배수 — 1.0 = 72 DPI. 1.5 ≈ 108 DPI (썸네일 + 가독성 균형).
_DEFAULT_ZOOM = 1.5

#: LibreOffice 변환 타임아웃 — 대형 PPT 대비 넉넉히, 단 무한 대기 X.
_SOFFICE_TIMEOUT_SECONDS = 60

#: 폰트 처리 모드 — env var `DOCUX_PREVIEW_FONT_MODE`로 토글.
#:   "original" (기본): **원본 폰트 우선**. 시스템에 그 폰트가 있으면 그대로 사용.
#:                      없을 때만 ``_PREVIEW_FONT_FAMILY``(Pretendard)로 fallback
#:                      (LibreOffice xcu substitution `Always=false`).
#:   "unified": 원본 폰트와 무관하게 ``_PREVIEW_FONT_FAMILY`` 강제 통일 (DOCX/PPTX/
#:              XLSX 정규화 + xcu `Always=true`). 포맷·스타일 일관 시각이 필요할 때.
_FONT_MODE = os.environ.get("DOCUX_PREVIEW_FONT_MODE", "original").lower()

#: 통일 폰트 또는 fallback 폰트 — env var `DOCUX_PREVIEW_FONT_FAMILY`로 override.
#: 디폴트 Pretendard — 한·영 통합 산세리프.
_PREVIEW_FONT_FAMILY = os.environ.get(
    "DOCUX_PREVIEW_FONT_FAMILY", "Pretendard"
)

#: "unified" 모드에서 LibreOffice substitution 대상 폰트.
_FONT_SUBSTITUTE_SOURCES = (
    "Calibri",
    "Cambria",
    "Arial",
    "Times New Roman",
    "Liberation Serif",
    "Liberation Sans",
    "Liberation Mono",
    "Apple SD Gothic Neo",
    "AppleGothic",
    "맑은 고딕",
    "Malgun Gothic",
    "바탕",
    "Batang",
    "굴림",
    "Gulim",
    "Dotum",
)

#: 격리된 LibreOffice user profile 디렉토리 캐시 — force(True/False) 별로 분리.
#: 두 종이 한 process 수명 동안 각각 1회 생성, 이후 재사용.
_LO_PROFILE_DIR: dict[str, Path] | None = None


def render_page_previews(
    path: Path,
    *,
    max_pages: int = _DEFAULT_MAX_PAGES,
    zoom: float = _DEFAULT_ZOOM,
    highlights: list[Highlight] | tuple[Highlight, ...] | None = None,
) -> dict[int, str]:
    """문서 페이지를 base64 PNG data URL로 렌더해 ``{page_no: data_url}``로 반환.

    페이지 번호는 1-base. 비PDF는 LibreOffice로 PDF 변환 후 동일 경로로 처리.
    ``highlights``가 주어지면 각 페이지에서 그 텍스트를 찾아 형광 annotation을
    추가해 픽맵에 반영한다(kind별 색: add→초록, delete→빨강, change→노랑).
    미지원 확장자·SDK 부재·변환 실패·렌더 실패는 빈 dict (보조 기능이라 본체 차단 X).
    한 페이지 실패는 다음 페이지 진행으로 폴백 (부분 결과 보존).
    """
    ext = path.suffix.lower()
    if ext not in _RENDERABLE_EXTS:
        return {}
    if ext == _PDF_EXT:
        return _render_pdf(
            path, max_pages=max_pages, zoom=zoom, highlights=highlights
        )
    if ext in _IMAGE_EXTS:
        # 스캔 이미지: pymupdf로 즉시 PDF 변환(LibreOffice 불요) 후 같은 경로로 렌더.
        converted = _image_to_temp_pdf(path)
        if converted is None:
            return {}
        try:
            return _render_pdf(
                converted, max_pages=max_pages, zoom=zoom, highlights=highlights
            )
        finally:
            converted.unlink(missing_ok=True)
    converted = _convert_to_pdf(path)
    if converted is None:
        return {}
    try:
        return _render_pdf(
            converted, max_pages=max_pages, zoom=zoom, highlights=highlights
        )
    finally:
        # 변환 산출물 즉시 정리 — 사내 자료 잔존 방지.
        try:
            converted.unlink(missing_ok=True)
        except OSError as exc:
            logger.warning("변환 PDF 정리 실패 (%s): %s", converted.name, exc)


def _render_pdf(
    path: Path,
    *,
    max_pages: int,
    zoom: float,
    highlights: list[Highlight] | tuple[Highlight, ...] | None = None,
) -> dict[int, str]:
    """PDF 한 파일을 페이지별 PNG data URL로 렌더 (pymupdf 단일 진입점).

    ``highlights``가 주어지면 페이지마다 `_apply_highlights`로 형광 annotation을
    그린 뒤 픽맵 렌더. pymupdf `get_pixmap()`은 기본으로 annotation을 같이 그리므로
    별도 합성 코드 없이 PNG에 색이 반영된다.
    """
    pymupdf = _load_pymupdf()
    if pymupdf is None:
        return {}
    try:
        doc = pymupdf.open(str(path))
    except Exception as exc:  # noqa: BLE001 — 손상 PDF·암호화 등은 빈 결과로 폴백.
        logger.warning("PDF 열기 실패 (%s): %s", path.name, exc)
        return {}
    try:
        previews: dict[int, str] = {}
        n_pages = min(len(doc), max_pages)
        matrix = pymupdf.Matrix(zoom, zoom)
        for i in range(n_pages):
            try:
                page = doc.load_page(i)
                if highlights:
                    _apply_highlights(page, highlights)
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                png_bytes = pix.tobytes("png")
                encoded = base64.b64encode(png_bytes).decode("ascii")
                previews[i + 1] = f"data:image/png;base64,{encoded}"
            except Exception as exc:  # noqa: BLE001 — 한 페이지 실패는 graceful skip.
                logger.warning(
                    "페이지 %d 렌더 실패 (%s): %s", i + 1, path.name, exc
                )
                continue
        return previews
    finally:
        doc.close()


def _tighten_punct_spacing(text: str) -> str:
    """괄호·구두점에 인접한 추출기 공백 인공물을 제거한다(단어 사이 공백은 보존).

    여는 괄호 뒤·닫는 괄호 앞·종결 구두점 앞의 공백만 제거한다. 예:
    '[ 런웨이 요금정책 ]' → '[런웨이 요금정책]' (PDF 원문과 일치해 search_for 매칭).
    """
    out = _RE_SPACE_AFTER_OPEN.sub(r"\1", text)
    out = _RE_SPACE_BEFORE_CLOSE.sub(r"\1", out)
    return out


def _prefix_candidates(text: str) -> list[str]:
    """긴 라인이 렌더 PDF에서 줄바꿈돼 full-line 검색이 빗나갈 때를 위한 점진적 단축 prefix.

    추출기 텍스트(docling 등)와 LibreOffice 변환 PDF의 텍스트 레이어는 줄바꿈
    위치가 달라, 라인 전체로는 `search_for`가 0건이 되기 쉽다. 라인 **시작**부터
    단어 경계로 점점 짧게 자른 부분문자열(연속)을 길이 내림차순으로 시도하면,
    최소한 라인 앞부분이라도 올바른 위치에 매칭돼 마커가 그려진다.

    단어 경계에서만 자른다(단어 중간 절단 → 매칭률 저하 방지). `_SEARCH_MIN_PREFIX`
    미만은 흔한 단어와 우연 매칭(false positive) 위험이 커 만들지 않는다.
    """
    words = text.split()
    out: list[str] = []
    for n in range(len(words) - 1, 0, -1):
        cand = " ".join(words[:n])
        if len(cand) < _SEARCH_MIN_PREFIX:
            break
        out.append(cand)
    return out


def _safe_search(page: object, text: str) -> list[object]:
    """page.search_for를 변형 후보 순서로 시도해 첫 매칭 rects를 반환(예외는 빈 결과).

    순서: 1) 원본 텍스트, 2) 구두점 공백 보정본, 3) 점진적 단축 prefix(단어 경계).
    추출기가 괄호·구두점 둘레에 끼운 공백 때문에 원본은 0건이어도 보정본이 매칭되거나,
    렌더 PDF의 줄바꿈으로 full-line이 빗나가도 짧은 prefix가 라인 앞부분에 매칭된다.
    """
    candidates = [text]
    tightened = _tighten_punct_spacing(text)
    if tightened != text:
        candidates.append(tightened)
    # full-line 후보 실패 시 점진적 단축 prefix — 보정본 기준(공백 인공물 제거된 형태).
    candidates.extend(_prefix_candidates(tightened))
    for cand in candidates:
        try:
            rects = page.search_for(cand)  # type: ignore[attr-defined]
        except Exception as exc:  # noqa: BLE001
            logger.debug("search_for 실패 (%r): %s", cand[:30], exc)
            return []
        if rects:
            return rects
    return []


def _apply_highlights(
    page: object, highlights: list[Highlight] | tuple[Highlight, ...]
) -> None:
    """페이지에 변경 텍스트별로 형광 사각형 annotation을 추가한다.

    알고리즘 (false positive 최소화):
      1. 라인 전체로 search_for(구두점 공백 보정 폴백 포함) → 라인 위치(rects) 확보.
      2. kind=="change" + changed_words 있음 → 각 변경 단어를 다시 검색하되,
         **라인 rect와 교차하는 단어 rect만** annotation. 페이지 다른 곳에
         있는 같은 단어는 무시 (예: 라인 A의 "50억" 변경, 페이지 다른 라인의
         "50억"은 색칠하지 않음).
      3. kind=="add"/"delete" 또는 changed_words 비어있음 → 라인 전체 색칠.
      4. 라인 매치 실패 → 그 entry 스킵 (조용한 폴백 X — 무관한 곳 색칠 방지).
    """
    for entry in highlights:
        line_text, kind, changed_words = entry
        normalized = _normalize_for_search(line_text)
        if not normalized:
            continue
        line_rects = _safe_search(page, normalized)
        if not line_rects:
            # 라인 매치 실패 — 페이지에 이 라인 없음(추출 텍스트 ↔ 렌더 PDF 레이어
            # 불일치 등). 무관한 곳 색칠 방지로 스킵하되, 마커 누락 진단용 로그를 남긴다.
            logger.debug("하이라이트 라인 매치 실패 — 스킵 (%r)", normalized[:40])
            continue
        color = _HIGHLIGHT_COLORS.get(kind, _HIGHLIGHT_COLORS["change"])

        if kind == "change" and changed_words:
            # 변경 단어 위치만 — 라인 rect와 교차하는 매치만 색칠.
            for word in changed_words:
                norm_word = _normalize_for_search(word)
                if not norm_word:
                    continue
                word_rects = _safe_search(page, norm_word)
                for wr in word_rects:
                    if not any(_rect_intersects(wr, lr) for lr in line_rects):
                        continue  # 라인 밖 매치는 false positive — 무시.
                    _annotate_rect(page, wr, color)
        else:
            # add/delete (또는 change words 없음): 라인 전체 색칠.
            for lr in line_rects:
                _annotate_rect(page, lr, color)


def _rect_intersects(a: object, b: object) -> bool:
    """pymupdf Rect 두 개가 교차하는지 (둘 다 axis-aligned)."""
    try:
        return not (
            a.x1 < b.x0  # type: ignore[attr-defined]
            or b.x1 < a.x0  # type: ignore[attr-defined]
            or a.y1 < b.y0  # type: ignore[attr-defined]
            or b.y1 < a.y0  # type: ignore[attr-defined]
        )
    except AttributeError:
        return False


def _annotate_rect(page: object, rect: object, color: tuple[float, float, float]) -> None:
    """rect에 형광 annotation을 건다 — 색·opacity 설정 + update."""
    try:
        annot = page.add_highlight_annot(rect)  # type: ignore[attr-defined]
        annot.set_colors(stroke=color)
        annot.set_opacity(_HIGHLIGHT_OPACITY)
        annot.update()
    except Exception as exc:  # noqa: BLE001
        logger.debug("annot 실패: %s", exc)


def _normalize_for_search(text: str | None) -> str | None:
    """검색 친화 형태로 정규화 — 마커·헤딩 프리픽스 제거, 길이 제한.

    None/빈 문자열·마커 라인·표 행·너무 짧은 텍스트는 None 반환 (검색 스킵).
    """
    if text is None:
        return None
    t = text.strip()
    if not t:
        return None
    # 마커 라인 스킵 — 추출기가 박은 [IMAGE …]·[TABLE …]·| 셀 |.
    if any(t.startswith(p) for p in _MARKER_PREFIXES):
        return None
    if t.startswith("|") and t.endswith("|"):
        return None
    # 마크다운 헤딩 프리픽스 제거 — PDF에는 # 없이 본문만 나옴.
    for prefix in ("### ", "## ", "# "):
        if t.startswith(prefix):
            t = t[len(prefix) :]
            break
    if len(t) < 2:
        return None
    if len(t) > _SEARCH_PREFIX_MAX:
        t = t[:_SEARCH_PREFIX_MAX]
    return t


def _image_to_temp_pdf(src: Path) -> Path | None:
    """스캔 이미지를 pymupdf로 1페이지 PDF(임시 파일)로 변환. 결과 Path(호출자 unlink).

    LibreOffice를 거치지 않고 pymupdf `convert_to_pdf`로 즉시 변환한다(빠름). 실패는
    None + WARNING(보조 기능이라 본체 차단 X).
    """
    pymupdf = _load_pymupdf()
    if pymupdf is None:
        return None
    try:
        doc = pymupdf.open(str(src))
        try:
            pdf_bytes = doc.convert_to_pdf()
        finally:
            doc.close()
        fd, tmp_name = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        tmp = Path(tmp_name)
        tmp.write_bytes(pdf_bytes)
        return tmp
    except Exception as exc:  # noqa: BLE001 — 손상 이미지 등은 None으로 폴백.
        logger.warning("이미지→PDF 변환 실패 (%s): %s", src.name, exc)
        return None


def _convert_to_pdf(src: Path) -> Path | None:
    """LibreOffice headless로 비PDF를 PDF로 변환. 결과는 임시 PDF Path (호출자 unlink).

    폰트 처리는 ``DOCUX_PREVIEW_FONT_MODE`` env var로 결정:
      - "original" (기본): 원본 폰트 그대로. 시스템에 없으면 LibreOffice
        metric-compatible 폴백(Carlito↔Calibri·Caladea↔Cambria·Liberation↔Arial/Times).
      - "unified": (1) DOCX/PPTX/XLSX는 사본에 모든 폰트를 ``_PREVIEW_FONT_FAMILY``로
        교체, (2) LibreOffice user profile에 substitution 룰 + 통일 폰트 강제.

    실패(soffice 부재·타임아웃·rc≠0·결과 PDF 부재)는 None + WARNING.
    """
    soffice = _find_soffice()
    if soffice is None:
        logger.warning(
            "LibreOffice(soffice) 미설치 — 비PDF 페이지 프리뷰 비활성 (%s). "
            "운영 환경에서는 컨테이너 이미지에 libreoffice를 포함하세요.",
            src.name,
        )
        return None
    unified = _FONT_MODE == "unified"
    # 양쪽 모드 모두 정규화 시도 — unified는 전부 강제, original은 시스템 부재 폰트만.
    normalized = _normalize_fonts(src)
    input_path = normalized if normalized is not None else src
    # xcu profile도 양쪽 모드 적용 — 미정규화 포맷(HTML/TXT/MD)의 추가 안전망.
    profile_dir = _ensure_lo_profile(force=unified)
    out_dir = Path(tempfile.mkdtemp(prefix="docux-pdf-"))
    try:
        cmd = [
            soffice,
            f"-env:UserInstallation=file://{profile_dir}",
        ]
        cmd += [
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(input_path),
        ]
        try:
            result = subprocess.run(  # noqa: S603 — 인자 리스트, shell=False, 신뢰 입력.
                cmd,
                capture_output=True,
                timeout=_SOFFICE_TIMEOUT_SECONDS,
                check=False,
            )
        except subprocess.TimeoutExpired:
            logger.warning(
                "LibreOffice 변환 타임아웃 (%s, %ds 상한)",
                src.name,
                _SOFFICE_TIMEOUT_SECONDS,
            )
            return None
        if result.returncode != 0:
            logger.warning(
                "LibreOffice 변환 실패 (%s): rc=%d", src.name, result.returncode
            )
            return None
        produced = out_dir / f"{input_path.stem}.pdf"
        if not produced.exists():
            logger.warning("변환 PDF 미발견 (%s)", src.name)
            return None
        # 결과 PDF만 호출자 lifecycle로 옮기고 임시 디렉토리는 즉시 정리.
        final_fd, final_name = tempfile.mkstemp(suffix=".pdf", prefix="docux-pdf-")
        os.close(final_fd)
        final_path = Path(final_name)
        shutil.move(str(produced), str(final_path))
        return final_path
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
        if normalized is not None and normalized.exists():
            try:
                normalized.unlink()
            except OSError:
                pass


def _normalize_fonts(src: Path) -> Path | None:
    """비PDF 파일 안 폰트를 _PREVIEW_FONT_FAMILY로 정규화한 임시 사본 반환.

    `_FONT_MODE`로 정책 결정:
      - "unified": **모든 폰트 강제 대체** (force_all=True).
      - "original": **원본 폰트가 시스템에 없을 때만** Pretendard로 대체 (fallback).

    DOCX/PPTX/XLSX만 처리(python-docx/pptx, openpyxl). 미지원 포맷·라이브러리
    부재·예외는 None — 호출자가 원본 그대로 변환한다 (graceful).
    """
    force_all = _FONT_MODE == "unified"
    ext = src.suffix.lower()
    try:
        if ext == ".docx":
            return _normalize_docx(src, force_all=force_all)
        if ext == ".pptx":
            return _normalize_pptx(src, force_all=force_all)
        if ext in {".xlsx", ".xlsm"}:
            return _normalize_xlsx(src, force_all=force_all)
    except Exception as exc:  # noqa: BLE001 — 정규화 실패는 원본 사용으로 graceful.
        logger.warning("폰트 정규화 실패 (%s): %s — 원본 그대로 변환", src.name, exc)
    return None


def _should_replace(font_name: str | None, *, force_all: bool) -> bool:
    """run-level 폰트 대체 여부.

    force_all=True → 무조건 대체.
    force_all=False → 명시된 폰트가 시스템에 없을 때만 대체.
                      None/빈 폰트는 LibreOffice가 디폴트 결정하므로 run-level
                      에서는 건드리지 않는다 (대신 style-level에서 처리).
    """
    if force_all:
        return True
    if not font_name:
        return False
    return not _font_available(font_name)


def _should_replace_style(font_name: str | None, *, force_all: bool) -> bool:
    """style-level(default) 폰트 대체 여부 — None도 적극 대체.

    style-level의 None은 LibreOffice가 임의 디폴트(보통 Calibri)를 선택해 시스템에
    없는 폰트로 폴백된다. 그래서 None이면 Pretendard로 명시해 LibreOffice 자체
    metric-compatible 폴백(Carlito 등)을 건너뛰게 한다.
    """
    if not font_name:
        return True
    if force_all:
        return True
    return not _font_available(font_name)


def _font_available(name: str) -> bool:
    """폰트가 시스템에 설치돼 있는지 — 정규화된 이름으로 매칭."""
    return _norm_font(name) in _system_font_families()


def _norm_font(name: str) -> str:
    return name.lower().replace(" ", "").replace("-", "").replace("_", "")


@lru_cache(maxsize=1)
def _system_font_families() -> frozenset[str]:
    """시스템 폰트 디렉토리에서 family 이름을 수집 (정규화된 lower-no-space)."""
    names: set[str] = set()
    dirs = [
        Path.home() / "Library" / "Fonts",
        Path("/Library/Fonts"),
        Path("/System/Library/Fonts"),
        Path("/System/Library/Fonts/Supplemental"),
        Path("/usr/share/fonts"),
        Path("/usr/local/share/fonts"),
    ]
    for d in dirs:
        if not d.exists():
            continue
        try:
            for f in d.rglob("*"):
                if f.suffix.lower() not in {".ttf", ".otf", ".ttc", ".otc"}:
                    continue
                stem = f.stem
                family = stem.split("-")[0].split("_")[0]
                names.add(_norm_font(stem))
                names.add(_norm_font(family))
        except OSError:
            continue
    return frozenset(names)


def _normalize_docx(src: Path, *, force_all: bool) -> Path | None:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return None
    # 1) 먼저 zipfile 레벨에서 theme1.xml의 majorFont/minorFont(헤딩·본문 디폴트)를
    #    Pretendard로 치환한 임시 사본 생성. python-docx는 theme XML을 노출 X.
    pre_patched = _patch_docx_theme(src, force_all=force_all)
    working_path = pre_patched if pre_patched is not None else src
    doc = Document(str(working_path))
    changed = pre_patched is not None
    # 1) docDefaults — 모든 style이 inherit하는 root 폰트 (Heading 등 포함).
    #    원본 모드에서도 force하는 이유: docDefaults는 보통 Calibri/Calibri Light
    #    (Microsoft 디폴트, macOS에 없음) → LibreOffice가 자동 Carlito 폴백.
    try:
        defaults = doc.styles.element.find(qn("w:docDefaults"))
        if defaults is not None:
            rpr_default = defaults.find(qn("w:rPrDefault"))
            if rpr_default is not None:
                rpr = rpr_default.find(qn("w:rPr"))
                if rpr is not None:
                    rfonts = rpr.find(qn("w:rFonts"))
                    if rfonts is None:
                        from docx.oxml import OxmlElement

                        rfonts = OxmlElement("w:rFonts")
                        rpr.insert(0, rfonts)
                    for attr in ("w:ascii", "w:eastAsia", "w:hAnsi", "w:cs"):
                        current = rfonts.get(qn(attr))
                        if _should_replace_style(current, force_all=force_all):
                            rfonts.set(qn(attr), _PREVIEW_FONT_FAMILY)
                            changed = True
    except Exception as exc:  # noqa: BLE001
        logger.debug("docDefaults 처리 실패: %s", exc)
    # 2) Style-level 폰트 — Heading 1, Normal 등 각 style의 rFonts 직접 조작.
    for style in doc.styles:
        try:
            style_element = getattr(style, "element", None)
            if style_element is None:
                continue
            rpr = style_element.find(qn("w:rPr"))
            if rpr is None:
                continue
            rfonts = rpr.find(qn("w:rFonts"))
            if rfonts is None:
                continue
            for attr in ("w:ascii", "w:eastAsia", "w:hAnsi", "w:cs"):
                current = rfonts.get(qn(attr))
                if _should_replace_style(current, force_all=force_all):
                    rfonts.set(qn(attr), _PREVIEW_FONT_FAMILY)
                    changed = True
        except Exception:  # noqa: BLE001
            continue
    # 3) Run-level 폰트 — 본문 paragraph runs.
    for para in doc.paragraphs:
        for run in para.runs:
            if _should_replace(run.font.name, force_all=force_all):
                run.font.name = _PREVIEW_FONT_FAMILY
                changed = True
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if _should_replace(run.font.name, force_all=force_all):
                            run.font.name = _PREVIEW_FONT_FAMILY
                            changed = True
    if not changed:
        if pre_patched is not None:
            pre_patched.unlink(missing_ok=True)
        return None
    result = _save_normalized(doc, src.suffix)
    if pre_patched is not None:
        pre_patched.unlink(missing_ok=True)
    return result


def _patch_docx_theme(src: Path, *, force_all: bool) -> Path | None:
    """DOCX의 theme1.xml에 정의된 majorFont/minorFont를 직접 교체한 임시 사본 반환.

    theme reference(`w:asciiTheme="minorHAnsi"`)는 styles/runs의 `w:ascii`보다
    원본 폰트 결정에 더 큰 영향. theme XML 자체를 패치해야 헤딩·본문 디폴트
    모두 `_PREVIEW_FONT_FAMILY`로 전환된다.
    """
    import zipfile

    try:
        from lxml import etree
    except ImportError:
        return None
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    fd, dst_name = tempfile.mkstemp(suffix=src.suffix, prefix="docux-theme-")
    os.close(fd)
    dst = Path(dst_name)
    any_change = False
    try:
        with zipfile.ZipFile(str(src)) as zin, zipfile.ZipFile(
            str(dst), "w", zipfile.ZIP_DEFLATED
        ) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item.startswith("word/theme/theme") and item.endswith(".xml"):
                    patched, did = _patch_theme_xml(
                        data, ns_a, etree, force_all=force_all
                    )
                    data = patched
                    any_change = any_change or did
                zout.writestr(item, data)
    except Exception as exc:  # noqa: BLE001
        logger.debug("theme 패치 실패: %s", exc)
        dst.unlink(missing_ok=True)
        return None
    if not any_change:
        dst.unlink(missing_ok=True)
        return None
    return dst


def _patch_theme_xml(
    xml_bytes: bytes, ns_a: str, etree: object, *, force_all: bool
) -> tuple[bytes, bool]:
    """theme XML의 majorFont/minorFont 안 모든 typeface를 _PREVIEW_FONT_FAMILY로.

    latin/ea/cs + ``<a:font script="...">``(스크립트별 폰트 — 한국어 "맑은 고딕",
    일본어 "ＭＳ 明朝" 등) 자식 모두 순회. 시스템에 없는 폰트를 Pretendard로 대체.
    """
    root = etree.fromstring(xml_bytes)  # type: ignore[attr-defined]
    changed = False
    for tag in ("majorFont", "minorFont"):
        for font_elem in root.iter(f"{{{ns_a}}}{tag}"):
            # latin/ea/cs (기본 자식)
            for child_tag in ("latin", "ea", "cs"):
                child = font_elem.find(f"{{{ns_a}}}{child_tag}")
                if child is None:
                    continue
                current = child.get("typeface")
                if _should_replace_style(current, force_all=force_all):
                    child.set("typeface", _PREVIEW_FONT_FAMILY)
                    changed = True
            # script별 자식 (Hang=한국어, Jpan=일본어, Hans/Hant=중국어 등)
            for script_font in font_elem.findall(f"{{{ns_a}}}font"):
                current = script_font.get("typeface")
                if _should_replace_style(current, force_all=force_all):
                    script_font.set("typeface", _PREVIEW_FONT_FAMILY)
                    changed = True
    return (
        etree.tostring(  # type: ignore[attr-defined]
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        ),
        changed,
    )


def _normalize_pptx(src: Path, *, force_all: bool) -> Path | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    prs = Presentation(str(src))
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _should_replace(run.font.name, force_all=force_all):
                        run.font.name = _PREVIEW_FONT_FAMILY
                        changed = True
    for master in prs.slide_masters:
        for shape in master.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _should_replace(run.font.name, force_all=force_all):
                        run.font.name = _PREVIEW_FONT_FAMILY
                        changed = True
    if not changed:
        return None
    return _save_normalized(prs, src.suffix)


def _normalize_xlsx(src: Path, *, force_all: bool) -> Path | None:
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Font
    except ImportError:
        return None
    wb = load_workbook(str(src))
    changed = False
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                existing = cell.font
                if not _should_replace(
                    existing.name if existing else None, force_all=force_all
                ):
                    continue
                cell.font = Font(
                    name=_PREVIEW_FONT_FAMILY,
                    size=existing.size if existing else None,
                    bold=existing.bold if existing else False,
                    italic=existing.italic if existing else False,
                    color=existing.color if existing else None,
                )
                changed = True
    if not changed:
        return None
    return _save_normalized(wb, src.suffix)


def _save_normalized(doc: object, suffix: str) -> Path:
    fd, name = tempfile.mkstemp(suffix=suffix, prefix="docux-norm-")
    os.close(fd)
    path = Path(name)
    doc.save(str(path))  # type: ignore[attr-defined]
    return path


def _ensure_lo_profile(*, force: bool) -> Path:
    """LibreOffice 격리 user profile 생성·재사용 + 폰트 substitution xcu 작성.

    `force=True`면 substitution `Always=true` — 원본 폰트가 시스템에 있어도 강제 대체.
    `force=False`면 `Always=false` — **원본 폰트 우선, 시스템에 없을 때만 fallback**.

    force별로 profile 디렉토리를 분리해 둘 다 캐시한다 (process 수명 1회).
    """
    global _LO_PROFILE_DIR
    if _LO_PROFILE_DIR is None:
        _LO_PROFILE_DIR = {}  # type: ignore[assignment]
    key = "force" if force else "fallback"
    cached = _LO_PROFILE_DIR.get(key) if isinstance(_LO_PROFILE_DIR, dict) else None  # type: ignore[union-attr]
    if cached is not None and cached.exists():
        return cached
    base = Path(tempfile.mkdtemp(prefix=f"docux-lo-profile-{key}-"))
    user_dir = base / "user"
    registry_file = user_dir / "registrymodifications.xcu"
    user_dir.mkdir(parents=True, exist_ok=True)
    registry_file.write_text(_build_registry_xml(force=force), encoding="utf-8")
    assert isinstance(_LO_PROFILE_DIR, dict)
    _LO_PROFILE_DIR[key] = base
    return base


def _build_registry_xml(*, force: bool) -> str:
    """폰트 substitution xcu 생성.

    force=True → Always=true (원본 무시, 강제 대체).
    force=False → Always=false (시스템에 원본 폰트 없을 때만 fallback).
    """
    always = "true" if force else "false"
    sub_nodes: list[str] = []
    for src in _FONT_SUBSTITUTE_SOURCES:
        sub_nodes.append(
            f'    <node oor:name="{src}" oor:op="replace">\n'
            f'      <prop oor:name="ReplaceFont" oor:op="fuse">'
            f"<value>{_PREVIEW_FONT_FAMILY}</value></prop>\n"
            f'      <prop oor:name="SubstituteFont" oor:op="fuse">'
            f"<value>{src}</value></prop>\n"
            f'      <prop oor:name="OnScreen" oor:op="fuse"><value>true</value></prop>\n'
            f'      <prop oor:name="Always" oor:op="fuse"><value>{always}</value></prop>\n'
            f"    </node>"
        )
    sub_block = "\n".join(sub_nodes)
    return f"""<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry"
           xmlns:xs="http://www.w3.org/2001/XMLSchema"
           xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <item oor:path="/org.openoffice.Office.Common/Font/Substitution">
    <prop oor:name="Replacement" oor:op="fuse"><value>true</value></prop>
  </item>
  <item oor:path="/org.openoffice.Office.Common/Font/Substitution/FontPairs">
{sub_block}
  </item>
</oor:items>
"""


def _find_soffice() -> str | None:
    """soffice / libreoffice 실행파일 위치 — PATH 우선, macOS 표준 경로 폴백."""
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    mac_app = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_app):
        return mac_app
    return None


def _load_pymupdf():  # type: ignore[no-untyped-def]
    """pymupdf 또는 구 이름 fitz를 lazy import. 둘 다 없으면 None + WARNING."""
    try:
        import pymupdf  # type: ignore[import-not-found]

        return pymupdf
    except ImportError:
        pass
    try:
        import fitz  # type: ignore[import-not-found]

        return fitz
    except ImportError:
        logger.warning(
            "pymupdf/fitz 미설치 — 페이지 프리뷰 비활성. "
            "`pip install pymupdf`로 활성화 가능."
        )
        return None


@lru_cache(maxsize=32)
def render_page_previews_highlighted_cached(
    path_str: str,
    mtime_ns: int,
    highlights_key: tuple[Highlight, ...],
) -> dict[int, str]:
    """(경로, mtime_ns, highlights) 키 캐시 — 같은 파일·같은 변경 집합이면 즉시 반환.

    `highlights_key`는 ``((text, kind, changed_words), ...)`` 튜플(hashable)이라 캐시
    키로 안전. 비어 있으면 일반 렌더와 동일 (annotation 없음).
    """
    highlights = list(highlights_key) if highlights_key else None
    return render_page_previews(Path(path_str), highlights=highlights)
