"""평문 폴백 파서 — 디지털·레이아웃 파서가 없거나 실패할 때의 마지막 보루.

각 포맷의 표준 라이브러리(python-docx, python-pptx, openpyxl)로 본문 텍스트를
뽑아 `MarkdownDoc`으로 환원한다. PDF는 PymupdfOcrParser가 별도로 처리하므로 여기서
다루지 않는다(중복 회피).

지원: .docx · .pptx · .xlsx · .xlsm · .html · .htm · .txt
"""

from __future__ import annotations

import logging
from pathlib import Path

from kms.adapters.ingestion.ir import MarkdownDoc

logger = logging.getLogger(__name__)

_SUPPORTED = {".docx", ".pptx", ".xlsx", ".xlsm", ".html", ".htm", ".txt", ".md"}


class PlainTextFallbackParser:
    """포맷별 표준 라이브러리로 평문을 뽑는 폴백 파서 (PDF 제외)."""

    name = "fallback"

    def is_available(self) -> bool:
        return True  # 표준 라이브러리 기반 — 항상 사용 가능 (의존성 미설치만 개별 체크).

    def supports(self, path: Path) -> bool:
        return path.suffix.lower() in _SUPPORTED

    def parse(self, path: Path) -> MarkdownDoc | None:
        ext = path.suffix.lower()
        try:
            if ext in {".txt", ".md"}:
                return _parse_txt(path)
            if ext == ".docx":
                return _parse_docx(path)
            if ext == ".pptx":
                return _parse_pptx(path)
            if ext in {".xlsx", ".xlsm"}:
                return _parse_xlsx(path)
            if ext in {".html", ".htm"}:
                return _parse_html(path)
        except Exception as exc:  # noqa: BLE001 — 폴백 실패는 None (코디네이터 처리).
            logger.warning("PlainText 폴백 실패 (%s): %s", ext, exc)
            return None
        return None


def _parse_txt(path: Path) -> MarkdownDoc:
    text = path.read_text(encoding="utf-8")
    return MarkdownDoc(markdown=text, page_map=[], image_blobs={})


def _parse_docx(path: Path) -> MarkdownDoc | None:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return None
    document = DocxDocument(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        # 헤딩 스타일 → 마크다운 헤더 변환 (markdown_chunker가 헤더 단위로 분할).
        if "heading 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style or "heading" in style:
            parts.append(f"### {text}")
        else:
            parts.append(text)
    # 표를 마크다운 표로 직렬화.
    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        ncols = max(len(r) for r in rows)
        parts.append(f"[TABLE r={len(rows)} c={ncols}]")
        for row in rows:
            padded = row + [""] * (ncols - len(row))
            parts.append("| " + " | ".join(padded) + " |")
    markdown = "\n\n".join(parts)
    return MarkdownDoc(markdown=markdown, page_map=[], image_blobs={})


def _parse_pptx(path: Path) -> MarkdownDoc | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    prs = Presentation(str(path))
    parts: list[str] = []
    page_map: list[tuple[int, int]] = []
    offset = 0
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"## 슬라이드 {idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                ncols = max((len(r) for r in rows), default=0)
                if rows and ncols:
                    slide_parts.append(f"[TABLE r={len(rows)} c={ncols}]")
                    for row in rows:
                        padded = row + [""] * (ncols - len(row))
                        slide_parts.append("| " + " | ".join(padded) + " |")
        slide_text = "\n\n".join(slide_parts) + "\n\n"
        page_map.append((offset, idx))
        parts.append(slide_text)
        offset += len(slide_text)
    return MarkdownDoc(markdown="".join(parts), page_map=page_map, image_blobs={})


def _parse_xlsx(path: Path) -> MarkdownDoc | None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # 빈 시트 스킵.
        nonempty = [r for r in rows if any(c not in (None, "") for c in r)]
        if not nonempty:
            continue
        parts.append(f"## {sheet_name}")
        ncols = max(len(r) for r in nonempty)
        parts.append(f"[TABLE r={len(nonempty)} c={ncols}]")
        for r in nonempty:
            cells = [("" if v is None else str(v)) for v in r]
            cells += [""] * (ncols - len(cells))
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("")
    wb.close()
    markdown = "\n\n".join(parts)
    return MarkdownDoc(markdown=markdown, page_map=[], image_blobs={})


def _parse_html(path: Path) -> MarkdownDoc | None:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    # 간단한 태그 제거 — Docling 미설치 환경의 최소 폴백.
    try:
        from html.parser import HTMLParser

        class _Stripper(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self.chunks: list[str] = []

            def handle_data(self, data: str) -> None:
                if data.strip():
                    self.chunks.append(data)

        stripper = _Stripper()
        stripper.feed(raw)
        text = "\n".join(stripper.chunks)
    except Exception:  # noqa: BLE001
        text = raw
    return MarkdownDoc(markdown=text, page_map=[], image_blobs={})
