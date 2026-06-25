"""추출 코디네이터 + 평문 폴백 단위 테스트.

기존 형식별 추출기(DocxExtractor·PptxExtractor·XlsxExtractor·TxtExtractor)는 제거되고
모든 포맷을 `PdfExtractor` 코디네이터가 3-파서 + 평문 폴백으로 처리한다.

Docling·Open-Parse는 lazy + 옵션 의존성이라 테스트 환경에서는 PlainTextFallback이
주 경로. 코디네이터의 우선순위·폴백·실패 처리는 Stub 파서로 결정론 검증.
"""

from pathlib import Path

import pandas as pd
import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from kms.adapters.ingestion.extract import (
    DoclingDigitalParser,
    OpenParseLayoutParser,
    PdfExtractor,
    PlainTextFallbackParser,
    PymupdfOcrParser,
    get_extractor,
)
from kms.adapters.ingestion.extract import registry as registry_module
from kms.adapters.ingestion.ir import MarkdownDoc


# ── 평문 폴백 파서 — 포맷별 ─────────────────────────────────────────────

def test_fallback_parses_txt(tmp_path: Path) -> None:
    path = tmp_path / "note.txt"
    path.write_text("줄1\n줄2", encoding="utf-8")

    result = PlainTextFallbackParser().parse(path)

    assert isinstance(result, MarkdownDoc)
    assert "줄1" in result.markdown
    assert "줄2" in result.markdown


def test_fallback_parses_docx(tmp_path: Path) -> None:
    path = tmp_path / "doc.docx"
    doc = DocxDocument()
    doc.add_heading("제목", level=1)
    doc.add_paragraph("본문 단락")
    doc.save(str(path))

    result = PlainTextFallbackParser().parse(path)

    assert isinstance(result, MarkdownDoc)
    assert "제목" in result.markdown
    assert "본문 단락" in result.markdown
    assert "# 제목" in result.markdown  # heading → 마크다운 헤더


def test_fallback_parses_pptx(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "슬라이드 제목"
    prs.save(str(path))

    result = PlainTextFallbackParser().parse(path)

    assert isinstance(result, MarkdownDoc)
    assert "슬라이드 제목" in result.markdown
    assert "## 슬라이드 1" in result.markdown


def test_fallback_parses_xlsx(tmp_path: Path) -> None:
    path = tmp_path / "data.xlsx"
    with pd.ExcelWriter(str(path)) as w:
        pd.DataFrame({"이름": ["가", "나"], "값": [1, 2]}).to_excel(
            w, sheet_name="데이터", index=False
        )

    result = PlainTextFallbackParser().parse(path)

    assert isinstance(result, MarkdownDoc)
    assert "데이터" in result.markdown
    assert "이름" in result.markdown
    assert "가" in result.markdown


def test_fallback_parses_html(tmp_path: Path) -> None:
    path = tmp_path / "page.html"
    path.write_text("<html><body><h1>제목</h1><p>본문</p></body></html>", encoding="utf-8")

    result = PlainTextFallbackParser().parse(path)

    assert isinstance(result, MarkdownDoc)
    assert "제목" in result.markdown
    assert "본문" in result.markdown


def test_fallback_supports_returns_true_for_known_exts(tmp_path: Path) -> None:
    parser = PlainTextFallbackParser()
    assert parser.supports(tmp_path / "a.txt")
    assert parser.supports(tmp_path / "a.docx")
    assert parser.supports(tmp_path / "a.pptx")
    assert parser.supports(tmp_path / "a.xlsx")
    assert parser.supports(tmp_path / "a.xlsm")
    assert parser.supports(tmp_path / "a.html")
    assert parser.supports(tmp_path / "a.htm")


# ── 코디네이터 (PdfExtractor) ─────────────────────────────────────────────


def test_coordinator_handles_all_formats(tmp_path: Path) -> None:
    """기본 코디네이터(파서 4개 + 폴백)로 다양한 포맷 추출."""
    txt = tmp_path / "note.txt"
    txt.write_text("간단 본문", encoding="utf-8")

    result = PdfExtractor().extract(txt)

    assert "간단 본문" in result.markdown


def test_coordinator_uses_first_available_parser() -> None:
    """우선순위대로 첫 성공 반환 (Stub 2개로 검증)."""
    class _Stub1:
        name = "digital"
        called = False
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            _Stub1.called = True
            return MarkdownDoc(markdown="from digital", page_map=[(0, 1)], image_blobs={})

    class _Stub2:
        name = "ocr"
        called = False
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            _Stub2.called = True
            return MarkdownDoc(markdown="from ocr", page_map=[(0, 1)], image_blobs={})

    extractor = PdfExtractor(parsers=[_Stub1(), _Stub2()])
    result = extractor.extract(Path("any.pdf"))

    assert result.markdown == "from digital"
    assert _Stub1.called is True
    assert _Stub2.called is False


def test_coordinator_falls_back_when_parser_returns_none() -> None:
    """1순위 None → 2순위 호출."""
    class _NoneP:
        name = "x"
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None: return None

    class _OkP:
        name = "y"
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            return MarkdownDoc(markdown="ok", page_map=[(0, 1)], image_blobs={})

    result = PdfExtractor(parsers=[_NoneP(), _OkP()]).extract(Path("any.pdf"))
    assert result.markdown == "ok"


def test_coordinator_skips_unavailable_parser() -> None:
    class _Missing:
        name = "digital"
        called = False
        def is_available(self) -> bool: return False
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            _Missing.called = True
            return MarkdownDoc(markdown="x", page_map=[(0, 1)], image_blobs={})

    class _OkP:
        name = "y"
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            return MarkdownDoc(markdown="ok", page_map=[(0, 1)], image_blobs={})

    PdfExtractor(parsers=[_Missing(), _OkP()]).extract(Path("any.pdf"))
    assert _Missing.called is False


def test_coordinator_skips_parser_that_does_not_support_format() -> None:
    """파서 supports()=False는 호출조차 X."""
    class _PdfOnly:
        name = "ocr"
        called = False
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return path.suffix == ".pdf"
        def parse(self, path: Path) -> MarkdownDoc | None:
            _PdfOnly.called = True
            return None

    class _Any:
        name = "any"
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return True
        def parse(self, path: Path) -> MarkdownDoc | None:
            return MarkdownDoc(markdown="ok", page_map=[(0, 1)], image_blobs={})

    PdfExtractor(parsers=[_PdfOnly(), _Any()]).extract(Path("doc.docx"))
    assert _PdfOnly.called is False  # PDF 전용은 docx에서 스킵


def test_coordinator_raises_when_no_parser_supports_extension() -> None:
    class _PdfOnly:
        name = "ocr"
        def is_available(self) -> bool: return True
        def supports(self, path: Path) -> bool: return path.suffix == ".pdf"
        def parse(self, path: Path) -> MarkdownDoc | None: return None

    with pytest.raises(RuntimeError, match="지원하는 파서 없음"):
        PdfExtractor(parsers=[_PdfOnly()]).extract(Path("note.txt"))


def test_default_parsers_include_four_types() -> None:
    """기본 코디네이터에 digital/layout/ocr/fallback 4종."""
    extractor = PdfExtractor()
    names = [p.name for p in extractor._parsers]  # type: ignore[attr-defined]
    assert names == ["digital", "layout", "ocr", "fallback"]


# ── 파서 supports() 카탈로그 ──────────────────────────────────────────────


@pytest.mark.parametrize(
    ("parser_cls", "ext", "expected"),
    [
        (DoclingDigitalParser, ".pdf", True),
        (DoclingDigitalParser, ".docx", True),
        (DoclingDigitalParser, ".pptx", True),
        (DoclingDigitalParser, ".xlsx", True),
        (DoclingDigitalParser, ".xlsm", True),
        (DoclingDigitalParser, ".html", True),
        (DoclingDigitalParser, ".txt", True),
        # 이미지도 docling OCR(IMAGE 포맷)로 처리 — 텍스트+좌표 추출 지원.
        (DoclingDigitalParser, ".png", True),
        (DoclingDigitalParser, ".jpg", True),
        (OpenParseLayoutParser, ".pdf", True),
        (OpenParseLayoutParser, ".docx", True),
        (OpenParseLayoutParser, ".html", True),
        (OpenParseLayoutParser, ".png", False),
        (PymupdfOcrParser, ".pdf", True),
        (PymupdfOcrParser, ".docx", False),  # OCR은 PDF 전용
        (PymupdfOcrParser, ".txt", False),
    ],
)
def test_parser_supports_matrix(
    parser_cls: type, ext: str, expected: bool
) -> None:
    parser = parser_cls()
    assert parser.supports(Path(f"a{ext}")) is expected


# ── 레지스트리 ──────────────────────────────────────────────────────────


@pytest.mark.parametrize(
    "filename",
    ["a.pdf", "a.docx", "a.pptx", "a.xlsx", "a.xlsm", "a.html", "a.htm",
     "a.txt", "a.md", "A.PDF"],
)
def test_registry_returns_coordinator_for_all_supported(filename: str) -> None:
    """모든 지원 확장자는 단일 코디네이터(PdfExtractor)로 라우팅된다."""
    extractor = get_extractor(Path(filename))
    assert isinstance(extractor, PdfExtractor)


def test_registry_raises_on_unsupported_format() -> None:
    with pytest.raises(registry_module.UnsupportedFormatError):
        get_extractor(Path("image.png"))
